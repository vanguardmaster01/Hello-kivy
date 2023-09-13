from kivy.core.window import Window
from kivy.lang import Builder
from kivy.uix.label import Label
from kivy.uix.screenmanager import ScreenManager, Screen
from kivy.clock import Clock
from kivy.uix.video import Video
from DbFuncs import db
import base64
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.videoplayer import VideoPlayer
from kivy.uix.image import Image
from pptx import Presentation
from config import utils
import os
from io import BytesIO
import subprocess

# import pyglet
# from pyglet.media import AVBinSource, StaticMemorySource, Player


class AdScreen(Screen):
    def __init__(self, **kwargs):
        super(AdScreen, self).__init__(**kwargs)
        Clock.schedule_once(self.retrieve_layout)

    def retrieve_layout(self, dt):
        ad = db.get_ad_row(35)
        if ad:
            if ad.type == 'MP4':
                temp_file = './img/ad.mp4'
                utils.write_to_file(ad.content, temp_file)

                videoPlayerLayout = VideoPlayerLayout(temp_file)
                print(f'self.manager{self.manager}')
                videoPlayerLayout.manager = self.manager
                videoPlayerLayout.bind(on_touch_up=videoPlayerLayout.on_video_touch_up)
                
                # Add the VideoPlayer widget to the BoxLayout
                self.add_widget(videoPlayerLayout)
            elif True:

                ppt_file = BytesIO(ad.content)

                # Save the blob data as a temporary .ppt file
                pptFileName = 'temp.ppt'
                pptxFileName = 'temp.pptx'
                temp_file = './' + pptFileName
                utils.write_to_file(ad.content, './' + pptFileName)

                # Convert .ppt file into .pptx file
                # Define the command as a list of arguments
                path = os.path.dirname( __file__ ) + '/'
                command = ['soffice', '--headless', '--convert-to', 'pptx', '--outdir', 
                           path, path + pptFileName]
                
                # Execute the command
                subprocess.run(command)


                presentation = Presentation(path + pptxFileName)
                slides = presentation.slides

                for slide in slides:
                    # Iterate through the shapes in the slide
                    for shape in slide.shapes:
                        # Check if the shape is an image
                        if shape.shape_type == 13:  # 13 corresponds to image shape type
                            # Extract the image
                            image = shape.image
                            image_bytes = image.blob
                            image_data = BytesIO(image_bytes)
                            
                            # Convert the BytesIO object to a base64-encoded string
                            base64_image = base64.b64encode(image_data.getvalue()).decode()
                            
                            # Create an Image widget with the base64-encoded string as the source
                            slide_image = Image(source=f"data:image/png;base64,{base64_image}", height=400)
                            slide_image.bind(on_touch_down=self.touch_screen)
                            self.add_widget(slide_image)
                
                os.remove('./' + pptFileName)
                os.remove('./' + pptxFileName)

        else:
            self.add_widget(Label(text='Ads not found', color=(1,0,0,1)))
        
    def touch_screen(self, instance, touch):
        if instance.collide_point(*touch.pos):
            self.manager.current = 'List'

class VideoPlayerLayout(BoxLayout):
    def __init__(self, temp_file, **kwargs):
        super(VideoPlayerLayout, self).__init__(**kwargs)
        self.manager = None
        self.temp_file = temp_file

        # Create a VideoPlayer widget
        self.player = VideoPlayer(source=temp_file, state='play',
                                  options={'eos': 'loop'})

        # Add the VideoPlayer widget to the BoxLayout
        self.add_widget(self.player)
    
    def on_video_touch_up(self, video, touch):
        # Handle the video player touch up event
        if video.collide_point(*touch.pos):
            self.manager.current = 'List'
            os.remove(self.temp_file)

